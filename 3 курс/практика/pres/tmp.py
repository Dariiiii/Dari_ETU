from pptx import Presentation

def find_lists_in_presentation(presentation_path):
    prs = Presentation(presentation_path)
    for slide_number, slide in enumerate(prs.slides,1):
        if slide_number == 8:
            print("Слайд №", slide_number)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    # Выводим уровень абзаца и текст, если они не пустые
                    print(paragraph.text)
            print()
    


# Путь к файлу презентации
presentation_path = 'Презентация_ВКР_Крыжановский.pptx'
find_lists_in_presentation(presentation_path)
